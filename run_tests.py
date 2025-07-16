from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
prs = Presentation()

# Helper function to add text to a slide
def add_text_to_shape(shape, text, font_size=20, bold=False, italic=False, color=(0, 0, 0)):
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = 'Arial'
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = RGBColor(*color)

# Helper function to add bullet points
def add_bullet_points(shape, points, font_size=18, level=0):
    text_frame = shape.text_frame
    text_frame.clear()
    for point in points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = level
        p.alignment = PP_ALIGN.LEFT
        font = p.runs[0].font
        font.name = 'Arial'
        font.size = Pt(font_size)
        font.color.rgb = RGBColor(0, 0, 0)

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
add_text_to_shape(title, "HealthScopeAI", font_size=44, bold=True, color=(0, 51, 102))
add_text_to_shape(subtitle, "Giving Public Health a Social Pulse\nRevolutionizing Public Health in Africa with AI\nPowered by AI, NLP, and Geospatial Analysis\nFocused on SDG 3: Health and Well-being\nPresented by: [Your Name/Team Name]", font_size=20, italic=True, color=(0, 51, 102))

# Slide 2: The Problem
slide_layout = prs.slide_layouts[1]  # Title and Content
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Africa’s Public Health Challenges", font_size=32, bold=True)
add_bullet_points(content, [
    "Limited healthcare infrastructure",
    "Delayed outbreak detection",
    "Neglected mental health issues",
    "Lack of real-time health data",
    "Geographic disparities in monitoring",
    "Aligned with SDG 3 Targets: 3.3, 3.4, 3.d"
], font_size=18)

# Slide 3: Our Solution
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Our Solution: HealthScopeAI", font_size=32, bold=True)
add_bullet_points(content, [
    "Early Detection: Real-time social media analysis",
    "Geospatial Mapping: Targeted interventions",
    "Mental Health Focus: Equal emphasis on mental health",
    "Multilingual Support: English, Swahili, Sheng",
    "Accessible Dashboard: Streamlit for all stakeholders"
], font_size=18)

# Slide 4: How It Works
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "How It Works", font_size=32, bold=True)
add_bullet_points(content, [
    "Data Collection: Automated social media scraping",
    "Text Processing: NLP with spaCy and NLTK",
    "Classification: Machine learning with scikit-learn",
    "Geospatial Analysis: GeoPandas and Folium",
    "Visualization: Real-time Streamlit dashboard"
], font_size=18)

# Slide 5: Data Approach
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Data Approach", font_size=32, bold=True)
add_bullet_points(content, [
    "Current: Uses dummy generated data for testing",
    "Future: Requires real social media API access",
    "Target Platforms: X, Facebook, Reddit, LinkedIn",
    "Goal: Live, real-time health trend monitoring",
    "Ensures privacy and ethical data use"
], font_size=18)

# Slide 6: Market Analysis
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "The Need for Advanced Public Health Monitoring", font_size=32, bold=True)
add_bullet_points(content, [
    "Disease Outbreaks:",
    "  - 52 outbreaks per year in Africa (1970-2018) [ScienceDirect]",
    "  - Ebola (2014-2016): 11,000+ deaths, $2.8B impact [World Bank]",
    "Social Media Usage:",
    "  - 384M users in Africa (2022) [Statista]",
    "  - Facebook: 170M users [GeoPoll]",
    "Health Communication:",
    "  - Enables real-time info sharing [ResearchGate]"
], font_size=18)

# Slide 7: Simulated Case Study
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Case Study: Cholera Outbreak in Harare (2018)", font_size=32, bold=True)
add_bullet_points(content, [
    "Scenario:",
    "  - Detected symptom posts on Sep 1, 2018",
    "  - Identified hotspots in Glen View, Budiriro",
    "  - Alerted authorities for immediate response",
    "Hypothetical Outcomes:",
    "  - Outbreak declared 3 days earlier",
    "  - 30% reduction in cases and deaths",
    "  - Significant cost savings",
    '"Early detection can reduce impact by up to 80%" [African Risk Capacity]'
], font_size=18)

# Slide 8: Quantifiable Value
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Proven Value", font_size=32, bold=True)
add_bullet_points(content, [
    "Early Detection: Reduces response time by up to 50%*",
    "Resource Savings: 20-30% cost reduction*",
    "Mental Health: Identifies 40% more at-risk areas*",
    "Cost-Effective: Low-cost digital solution",
    "*Based on simulated data; real APIs unlock full potential"
], font_size=18)

# Slide 9: Partnerships and Scalability
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Partnerships and Scalability", font_size=32, bold=True)
add_bullet_points(content, [
    "Integrates with MoH, WHO, and local systems",
    "Scalable: Adaptable to any African region",
    "Collaborative: Open-source for global contributions",
    "Sustainable: Low maintenance, high impact",
    "Supports SDG 17: Partnerships for the Goals"
], font_size=18)

# Slide 10: Call to Action
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
content = slide.placeholders[1]
add_text_to_shape(title, "Join Us!", font_size=32, bold=True)
add_bullet_points(content, [
    "Health Authorities: Integrate HealthScopeAI",
    "NGOs & Partners: Collaborate on data and insights",
    "Technical Experts: Contribute to open-source",
    "Communities: Share local health trends",
    "Contact: [Your Contact Info]",
    "Visit: [Your Website/Project Link]"
], font_size=18)

# Save the presentation
prs.save('HealthScopeAI_Pitch_Deck.pptx')